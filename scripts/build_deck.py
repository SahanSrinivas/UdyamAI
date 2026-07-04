#!/usr/bin/env python3
"""Fill the IDBI Innovate template with UdyamAI content on slides 3-13.

Opens the official IDBI template, preserves every existing shape (header, footer,
title text box), and layers modern content into the empty area below the title.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


def E(v):
    return Emu(int(v))


ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "Prototype Submission Deck _ IDBI Innovate.pptx"
OUT = ROOT / "UdyamAI - IDBI Innovate Submission.pptx"
SHOTS = ROOT / "deck-screenshots"

# ─── palette ───────────────────────────────────────────────────────────
NAVY = RGBColor(10, 22, 40)
NAVY_SOFT = RGBColor(30, 50, 90)
INK = RGBColor(20, 25, 34)
INK_2 = RGBColor(40, 48, 62)
LIME = RGBColor(204, 255, 94)
LIME_DIM = RGBColor(148, 200, 60)
SAFFRON = RGBColor(255, 153, 51)
GREEN = RGBColor(19, 136, 8)
WHITE = RGBColor(255, 255, 255)
CREAM = RGBColor(250, 247, 241)
SAND = RGBColor(237, 229, 214)
MIST = RGBColor(216, 224, 236)
LINE = RGBColor(224, 226, 232)
LINE_2 = RGBColor(200, 206, 216)
TEXT = RGBColor(20, 25, 34)
MUTED = RGBColor(96, 105, 116)
MUTED_2 = RGBColor(140, 148, 160)
RED = RGBColor(220, 55, 55)
BLUE = RGBColor(37, 99, 235)
GOLD = RGBColor(207, 165, 79)
AWS_ORANGE = RGBColor(255, 153, 0)
AWS_DARK = RGBColor(35, 47, 62)
IDBI_MAROON = RGBColor(120, 20, 30)

# ─── typography ────────────────────────────────────────────────────────
SERIF = "Georgia"
SANS = "Helvetica Neue"
MONO = "Menlo"

# Template geometry
CONTENT_LEFT = Inches(0.35)
CONTENT_TOP = Inches(1.55)
CONTENT_RIGHT = Inches(9.65)
CONTENT_BOTTOM = Inches(5.35)
CONTENT_W = CONTENT_RIGHT - CONTENT_LEFT
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP


# ─── primitives ────────────────────────────────────────────────────────

def rect(slide, left, top, width, height, fill_color, line_color=None, radius=False,
         radius_amt=0.12):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, E(left), E(top), E(width), E(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius_amt
    return shape


def oval(slide, left, top, size, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(left), E(top), E(size), E(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, left, top, width, height, text, size=11, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=SANS, italic=False,
        line_spacing=None):
    box = slide.shapes.add_textbox(E(left), E(top), E(width), E(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = p if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        if line_spacing:
            para.line_spacing = line_spacing
        for run in para.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return box


def arrow(slide, x1, y1, x2, y2, color=NAVY_SOFT, width=1.25, with_head=True):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, E(x1), E(y1),
                                       E(x2), E(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if with_head:
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "sm", "len": "sm"})
        ln.append(tail)
    return conn


def hline(slide, left, top, width, color=LINE, weight=0.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, E(left), E(top),
                                       E(left + width), E(top))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ─── logo (reconstructed in pptx shapes) ───────────────────────────────

def draw_logo(slide, left, top, size=Inches(0.7)):
    """UdyamAI mark: lime rounded square with a serif U inscribed."""
    rect(slide, left, top, size, size, LIME, radius=True, radius_amt=0.24)
    # Inner shadow ring (subtle, gives depth)
    inset = size / 12
    rect(slide, left + inset, top + inset, size - inset * 2, size - inset * 2,
         LIME, radius=True, radius_amt=0.22)
    # Big serif U
    txt(slide, left, top, size, size, "U",
        size=int(size / 914400 * 42), bold=True, color=INK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SERIF)


def draw_wordmark(slide, left, top, size=Inches(0.6)):
    """Full lockup: logo mark + Udyam + AI (italic serif)."""
    draw_logo(slide, left, top, size)
    gap = size + size / 6
    w = Inches(2.8)
    # "Udyam" in sans + "AI" in italic serif
    box = slide.shapes.add_textbox(E(left + gap), E(top), E(w), E(size))
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = ""
    r1 = p.add_run()
    r1.text = "Udyam"
    r1.font.name = SANS
    r1.font.size = Pt(int(size / 914400 * 28))
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = "AI"
    r2.font.name = SERIF
    r2.font.size = Pt(int(size / 914400 * 28))
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = LIME_DIM


# ─── higher-level components ───────────────────────────────────────────

def style_cell(cell, text_str, bold=False, size=9, color=TEXT, fill=None,
               align=PP_ALIGN.LEFT, font=SANS):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(7)
    cell.margin_right = Pt(7)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    p = tf.paragraphs[0]
    p.text = text_str
    p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def table(slide, left, top, width, height, headers, rows, col_widths=None,
          header_fill=NAVY, alt_fill=CREAM, header_size=9, row_size=8):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, E(left), E(top), E(width),
                                        E(height))
    t = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = E(w)
    for j, h in enumerate(headers):
        style_cell(t.cell(0, j), h, bold=True, size=header_size, color=WHITE,
                   fill=header_fill, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        bg = alt_fill if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            is_total = isinstance(val, str) and (val.startswith("Total") or
                                                 (j == 0 and val.lower().startswith("total")))
            bold = (j == 0) or is_total
            fill = alt_fill if is_total else bg
            color = NAVY if is_total else TEXT
            style_cell(t.cell(i, j), str(val), bold=bold, size=row_size,
                       color=color, fill=fill)
    return t


def stat_pill(slide, left, top, label, value, accent=SAFFRON, width=Inches(2.2),
              height=Inches(0.75), value_size=16, label_size=8):
    rect(slide, left, top, width, height, WHITE, LINE, radius=True)
    rect(slide, left, top, Inches(0.09), height, accent)
    txt(slide, left + Inches(0.22), top + Inches(0.08),
        width - Inches(0.28), Inches(0.32), value,
        size=value_size, bold=True, color=NAVY)
    txt(slide, left + Inches(0.22), top + height - Inches(0.3),
        width - Inches(0.28), Inches(0.24), label.upper(),
        size=label_size, bold=True, color=MUTED)


def callout(slide, left, top, width, text_str, accent=LIME, height=Inches(0.55)):
    rect(slide, left, top, width, height, NAVY, radius=True)
    rect(slide, left, top, Inches(0.08), height, accent)
    txt(slide, left + Inches(0.22), top, width - Inches(0.3), height, text_str,
        size=10.5, bold=True, color=LIME,
        anchor=MSO_ANCHOR.MIDDLE)


def section_label(slide, left, top, width, label, color=MUTED):
    txt(slide, left, top, width, Inches(0.22), label.upper(),
        size=8, bold=True, color=color)
    hline(slide, left, top + Inches(0.24), Inches(0.35), color=color, weight=1.5)


def aws_badge(slide, left, top, width, height, service, detail):
    """AWS-styled service card — orange top strip + navy service name + muted detail."""
    rect(slide, left, top, width, height, WHITE, LINE_2, radius=True)
    rect(slide, left, top, width, Inches(0.07), AWS_ORANGE, radius=True)
    # AWS wordmark micro-tag
    txt(slide, left + Inches(0.12), top + Inches(0.12),
        Inches(0.55), Inches(0.18),
        "AWS", size=7, bold=True, color=AWS_ORANGE)
    txt(slide, left + Inches(0.12), top + Inches(0.32),
        width - Inches(0.2), Inches(0.24),
        service, size=9, bold=True, color=AWS_DARK)
    if detail:
        txt(slide, left + Inches(0.12), top + Inches(0.56),
            width - Inches(0.2), Inches(0.24),
            detail, size=7.5, color=MUTED)


# ─── SLIDE 2 · Brief about the idea (hero pitch + logo) ────────────────

def slide_brief_idea(slide):
    """Slide 2 already has 'Brief about the idea' title and subtitle.
    Layer the full pitch statement + logo + impact stats below.
    Available content area: y ~ 1.85 to 5.35."""

    S_TOP = Inches(1.85)
    S_BOTTOM = Inches(5.35)

    # Small logo mark top-right (aligned with subtitle area)
    draw_logo(slide, Inches(9.0), Inches(1.28), size=Inches(0.44))

    # Section label + rule
    section_label(slide, CONTENT_LEFT, S_TOP, Inches(3), "The pitch",
                  color=SAFFRON)

    # Hero line — Georgia serif, bold, editorial
    hero_top = S_TOP + Inches(0.32)
    hero_text = (
        "Banks reject 70%+ of MSME applications because NTC / NTB enterprises "
        "can't prove creditworthiness with documents they don't have."
    )
    txt(slide, CONTENT_LEFT, hero_top, CONTENT_W, Inches(0.7),
        hero_text, size=14.5, bold=True, color=INK, font=SERIF,
        line_spacing=1.18)

    # Sub-pitch — clean sans, muted
    sub_top = hero_top + Inches(0.78)
    sub_text = (
        "UdyamAI is the unified assessment framework on AA + GST + UPI + EPFO — "
        "a living Health Card for the MSME, a pre-qualified lead engine for the bank, "
        "and an audit-first ML layer RBI can approve. We don't score you at apply time. "
        "We score you continuously — so rejection becomes fixable, and portfolio quality "
        "improves before the loan is booked."
    )
    txt(slide, CONTENT_LEFT, sub_top, CONTENT_W, Inches(1.0),
        sub_text, size=10, color=INK_2, line_spacing=1.28)

    # 3 impact stat pills
    stats_top = Inches(3.95)
    pw = (CONTENT_W - Inches(0.4)) / 3
    stat_pill(slide, CONTENT_LEFT, stats_top, "Current rejection rate",
              "70%+", RED, pw, Inches(0.72), value_size=20, label_size=7.5)
    stat_pill(slide, CONTENT_LEFT + pw + Inches(0.2), stats_top,
              "Credit-invisible enterprises", "~14 M", SAFFRON, pw, Inches(0.72),
              value_size=20, label_size=7.5)
    stat_pill(slide, CONTENT_LEFT + 2 * (pw + Inches(0.2)), stats_top,
              "MSME credit gap · India", "₹25 L Cr", GREEN, pw, Inches(0.72),
              value_size=20, label_size=7.5)

    # Comparison strip at bottom — 2 cards
    cmp_top = Inches(4.75)
    cmp_h = Inches(0.55)
    lw = (CONTENT_W - Inches(0.15)) / 2

    rect(slide, CONTENT_LEFT, cmp_top, lw, cmp_h, CREAM, LINE, radius=True)
    txt(slide, CONTENT_LEFT + Inches(0.16), cmp_top + Inches(0.06),
        Inches(1.6), Inches(0.2),
        "EXISTING PLAYERS", size=7, bold=True, color=MUTED)
    txt(slide, CONTENT_LEFT + Inches(0.16), cmp_top + Inches(0.26),
        lw - Inches(0.25), Inches(0.28),
        "Score you once, at apply time.  \"Try again later.\"",
        size=8.5, color=MUTED_2, italic=True)

    rx = CONTENT_LEFT + lw + Inches(0.15)
    rect(slide, rx, cmp_top, lw, cmp_h, NAVY, radius=True)
    rect(slide, rx, cmp_top, Inches(0.08), cmp_h, LIME)
    txt(slide, rx + Inches(0.2), cmp_top + Inches(0.06),
        Inches(1.4), Inches(0.2),
        "UDYAMAI", size=7, bold=True, color=LIME)
    txt(slide, rx + Inches(0.2), cmp_top + Inches(0.26),
        lw - Inches(0.28), Inches(0.28),
        "Continuous · fixable · pre-qualified across 64 lenders · EN · HI · TE",
        size=8.5, bold=True, color=WHITE)


# ─── SLIDE 3 · Opportunities (how different · how solve · USP) ─────────

def slide_opportunities(slide):
    # Small logo mark top-right for brand continuity
    draw_logo(slide, Inches(9.0), CONTENT_TOP - Inches(0.55), size=Inches(0.44))

    # Left column — comparison table (answers "how different")
    section_label(slide, CONTENT_LEFT, CONTENT_TOP, Inches(3),
                  "How we're different", color=SAFFRON)

    tbl_left = CONTENT_LEFT
    tbl_top = CONTENT_TOP + Inches(0.32)
    tbl_w = Inches(5.4)
    tbl_h = Inches(2.55)

    headers = ["Dimension", "Existing players", "UdyamAI"]
    rows = [
        ["Audience", "Sell scoring to banks", "Sell to the borrower"],
        ["Score cadence", "One-shot at apply time", "Continuous · monthly"],
        ["On rejection", "\"Try again later\"", "\"Fix this exact number\""],
        ["Lender coverage", "One at a time", "64 lenders via ULI"],
        ["Language", "English-only", "EN · HI · TE"],
        ["ML transparency", "Black-box XGBoost", "Auditable LR (β visible)"],
    ]
    table(slide, tbl_left, tbl_top, tbl_w, tbl_h, headers, rows,
          col_widths=[Inches(1.0), Inches(2.2), Inches(2.2)],
          header_size=9, row_size=8.5)

    # Right column — solutions ("how it solves the problem")
    rc_left = CONTENT_LEFT + tbl_w + Inches(0.25)
    rc_w = CONTENT_RIGHT - rc_left
    section_label(slide, rc_left, CONTENT_TOP, rc_w, "How it solves it",
                  color=GREEN)

    solve = [
        ("Aggregation of AA + GST + UPI + EPFO",
         "One consent · four data rails unified into one score"),
        ("Quantified nudges",
         "\"+34 points if you file GSTR-3B on time\" — actionable, numeric"),
        ("Cross-lender OCEN",
         "Portable score across 64 ULI-onboarded banks · no lender lock-in"),
        ("Vernacular Gemini AI",
         "Every score change explained in EN · HI · TE"),
    ]
    ry = tbl_top
    for i, (title, sub) in enumerate(solve):
        accent = [SAFFRON, GREEN, BLUE, NAVY][i]
        rect(slide, rc_left, ry, rc_w, Inches(0.58), WHITE, LINE, radius=True)
        rect(slide, rc_left, ry, Inches(0.08), Inches(0.58), accent)
        # number circle
        oval(slide, rc_left + Inches(0.18), ry + Inches(0.15), Inches(0.28),
             accent)
        txt(slide, rc_left + Inches(0.18), ry + Inches(0.15),
            Inches(0.28), Inches(0.28), str(i + 1),
            size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, rc_left + Inches(0.56), ry + Inches(0.08),
            rc_w - Inches(0.65), Inches(0.24), title,
            size=9.5, bold=True, color=NAVY)
        txt(slide, rc_left + Inches(0.56), ry + Inches(0.3),
            rc_w - Inches(0.65), Inches(0.28), sub,
            size=8, color=MUTED)
        ry += Inches(0.63)

    # USP callout across the bottom
    callout(slide, CONTENT_LEFT, CONTENT_BOTTOM - Inches(0.55), CONTENT_W,
            "USP · The only credit dashboard where your MSME score improves before you apply.")


# ─── SLIDE 4 · Features · dual persona (MSME + IDBI banker) ────────────

def slide_features(slide):
    """Two big columns — one for the MSME borrower, one for the IDBI banker."""

    cw = (CONTENT_W - Inches(0.2)) / 2
    ch = CONTENT_H - Inches(0.05)

    # ─── LEFT: For the MSME ────────────────────────────────
    x = CONTENT_LEFT
    y = CONTENT_TOP + Inches(0.02)
    rect(slide, x, y, cw, ch, WHITE, LINE, radius=True)
    rect(slide, x, y, cw, Inches(0.62), NAVY, radius=True)
    # persona badge
    rect(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.28), Inches(0.28),
         LIME, radius=True)
    txt(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.28), Inches(0.28),
        "M", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=SERIF)
    txt(slide, x + Inches(0.55), y + Inches(0.12),
        cw - Inches(0.6), Inches(0.24),
        "FOR THE MSME OWNER", size=8, bold=True, color=LIME)
    txt(slide, x + Inches(0.55), y + Inches(0.32),
        cw - Inches(0.6), Inches(0.28),
        "Living Health Card", size=14, bold=True, color=WHITE)

    msme_feats = [
        ("4D Health Score",
         "Revenue · Compliance · Counterparty · Growth — 0-1000"),
        ("Quantified nudges",
         "\"+34 points if you file GSTR-3B on time\""),
        ("Vernacular explainer",
         "Gemini AI in English · Hindi · Telugu"),
        ("Continuous refresh",
         "Score updates the day your GSTR-3B is filed"),
        ("NTC / NTB / Thin-file support",
         "Works without CIBIL, without audited books"),
    ]
    iy = y + Inches(0.85)
    for title, sub in msme_feats:
        oval(slide, x + Inches(0.22), iy + Inches(0.12), Inches(0.1), LIME)
        txt(slide, x + Inches(0.42), iy, cw - Inches(0.55), Inches(0.22),
            title, size=9.5, bold=True, color=NAVY)
        txt(slide, x + Inches(0.42), iy + Inches(0.22),
            cw - Inches(0.55), Inches(0.22),
            sub, size=8, color=MUTED)
        iy += Inches(0.5)

    # ─── RIGHT: For IDBI Bank ──────────────────────────────
    x = CONTENT_LEFT + cw + Inches(0.2)
    y = CONTENT_TOP + Inches(0.02)
    rect(slide, x, y, cw, ch, WHITE, LINE, radius=True)
    rect(slide, x, y, cw, Inches(0.62), IDBI_MAROON, radius=True)
    rect(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.28), Inches(0.28),
         SAFFRON, radius=True)
    txt(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.28), Inches(0.28),
        "I", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=SERIF)
    txt(slide, x + Inches(0.55), y + Inches(0.12),
        cw - Inches(0.6), Inches(0.24),
        "FOR IDBI BANK · MSME DESK", size=8, bold=True, color=SAFFRON)
    txt(slide, x + Inches(0.55), y + Inches(0.32),
        cw - Inches(0.6), Inches(0.28),
        "Pre-qualified pipeline", size=14, bold=True, color=WHITE)

    idbi_feats = [
        ("Ranked lead pipeline",
         "LR-calibrated · approval confidence per MSME"),
        ("Rejection driver analytics",
         "\"42% of rejections fail Compliance\" — actionable"),
        ("Portfolio diversification signal",
         "Sector · geography · cohort concentration alerts"),
        ("One-click JSON export",
         "Straight to your underwriting desk"),
        ("Post-sanction monitoring",
         "OCEN status webhook · early-warning at -80 pts"),
    ]
    iy = y + Inches(0.85)
    for title, sub in idbi_feats:
        oval(slide, x + Inches(0.22), iy + Inches(0.12), Inches(0.1), SAFFRON)
        txt(slide, x + Inches(0.42), iy, cw - Inches(0.55), Inches(0.22),
            title, size=9.5, bold=True, color=IDBI_MAROON)
        txt(slide, x + Inches(0.42), iy + Inches(0.22),
            cw - Inches(0.55), Inches(0.22),
            sub, size=8, color=MUTED)
        iy += Inches(0.5)


# ─── SLIDE 5 · Process Flow ────────────────────────────────────────────

def slide_flowchart(slide):
    cx = (CONTENT_LEFT + CONTENT_RIGHT) / 2
    y = CONTENT_TOP + Inches(0.05)

    def hbox(text_str, y_pos, fill, tc=WHITE, w=Inches(2.2), h=Inches(0.4)):
        x = cx - w / 2
        rect(slide, x, y_pos, w, h, fill, radius=True)
        txt(slide, x, y_pos, w, h, text_str, size=9, bold=True, color=tc,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return y_pos + h

    def arr(y_from, drop=Inches(0.13)):
        arrow(slide, cx, y_from, cx, y_from + drop)
        return y_from + drop

    y = hbox("MSME enters GSTIN", y, SAFFRON, WHITE)
    y = arr(y)
    y = hbox("Account Aggregator · consent flow", y, NAVY, WHITE, w=Inches(2.8))
    y = arr(y)

    # 5 data sources
    sources = ["AA · Bank", "GSTN", "UPI", "EPFO", "Bureau*"]
    sw = Inches(1.2)
    gap = Inches(0.08)
    total = sw * len(sources) + gap * (len(sources) - 1)
    sx = cx - total / 2
    src_y = y
    for i, s in enumerate(sources):
        x = sx + i * (sw + gap)
        rect(slide, x, src_y, sw, Inches(0.36), CREAM, NAVY_SOFT, radius=True)
        txt(slide, x, src_y, sw, Inches(0.36), s, size=8, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        arrow(slide, x + sw / 2, src_y + Inches(0.36), cx, src_y + Inches(0.56))

    y = src_y + Inches(0.58)
    y = hbox("UdyamAI Score Engine · 4 sub-scores → 1000", y, GREEN, WHITE,
             w=Inches(3.3), h=Inches(0.42))
    y = arr(y, Inches(0.11))

    outs = [("Nudge engine", SAFFRON), ("Gemini Explainer", BLUE),
            ("LR Calibrator", NAVY)]
    ow = Inches(1.85)
    ogap = Inches(0.12)
    otot = ow * 3 + ogap * 2
    ox = cx - otot / 2
    out_y = y
    for i, (label, color) in enumerate(outs):
        x = ox + i * (ow + ogap)
        tc = TEXT if color == SAFFRON else WHITE
        rect(slide, x, out_y, ow, Inches(0.4), color, radius=True)
        txt(slide, x, out_y, ow, Inches(0.4), label, size=9, bold=True, color=tc,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        arrow(slide, x + ow / 2, out_y + Inches(0.4), cx, out_y + Inches(0.6))

    y = out_y + Inches(0.6)
    stages = [("Pre-qualified quotes", LIME, TEXT),
              ("OCEN LA → Lender", NAVY, WHITE),
              ("Sanction · 24h", GREEN, WHITE)]
    sw2 = Inches(2.3)
    sgap = Inches(0.25)
    stot = sw2 * 3 + sgap * 2
    sxp = cx - stot / 2
    for i, (label, fill, tc) in enumerate(stages):
        x = sxp + i * (sw2 + sgap)
        rect(slide, x, y, sw2, Inches(0.44), fill, radius=True)
        txt(slide, x, y, sw2, Inches(0.44), label, size=9.5, bold=True, color=tc,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            arrow(slide, x + sw2, y + Inches(0.22),
                  x + sw2 + sgap, y + Inches(0.22), NAVY_SOFT)


# ─── screenshots (6, 7, 10) ─────────────────────────────────────────────

def slide_screenshots(slide):
    labels = [
        "01 · Landing hero",
        "02 · GSTIN validator",
        "03 · Dashboard overview",
        "04 · Nudge panel",
        "05 · Loan quotes + LR",
        "06 · OCEN drawer",
    ]
    cols, rows_ = 3, 2
    gw = (CONTENT_W - Inches(0.2)) / cols
    gh = (CONTENT_H - Inches(0.1)) / rows_
    gap = Inches(0.1)
    for idx in range(6):
        r, c = divmod(idx, cols)
        x = CONTENT_LEFT + c * (gw + gap)
        y = CONTENT_TOP + r * (gh + Inches(0.05))
        rect(slide, x, y, gw, gh, WHITE, LINE, radius=True)
        shot = SHOTS / f"{idx + 1:02d}.png"
        img_top = y + Inches(0.08)
        img_h = gh - Inches(0.42)
        if shot.exists():
            slide.shapes.add_picture(str(shot), E(x + Inches(0.08)), E(img_top),
                                     E(gw - Inches(0.16)), E(img_h))
        else:
            txt(slide, x, img_top, gw, img_h,
                "[ Screenshot pending ]", size=9, color=MUTED,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, x, y + gh - Inches(0.3), gw, Inches(0.3), CREAM)
        txt(slide, x + Inches(0.12), y + gh - Inches(0.25),
            gw - Inches(0.16), Inches(0.22), labels[idx],
            size=8, bold=True, color=NAVY)


# ─── SLIDE 7 · Architecture — AWS-native (per IDBI toolkit) ────────────

def slide_architecture(slide):
    """AWS-native architecture. Left = layered app stack with AWS badges,
    right = AWS foundation services + regulated rails at the bottom."""

    # "Built on AWS" ribbon at the very top-right
    ribbon_w = Inches(1.7)
    rect(slide, CONTENT_RIGHT - ribbon_w, CONTENT_TOP - Inches(0.05),
         ribbon_w, Inches(0.32), AWS_DARK, radius=True)
    rect(slide, CONTENT_RIGHT - ribbon_w, CONTENT_TOP - Inches(0.05),
         Inches(0.08), Inches(0.32), AWS_ORANGE)
    txt(slide, CONTENT_RIGHT - ribbon_w + Inches(0.15),
        CONTENT_TOP - Inches(0.03), ribbon_w - Inches(0.2), Inches(0.28),
        "Built on AWS · ACC",
        size=8, bold=True, color=AWS_ORANGE,
        anchor=MSO_ANCHOR.MIDDLE)

    # ─── LEFT column: main app stack ───────────────
    main_left = CONTENT_LEFT
    main_w = Inches(6.0)
    top = CONTENT_TOP + Inches(0.4)

    layers = [
        ("CLIENT", SAFFRON, [
            ("Mobile · iOS / Android", "React Native"),
            ("Web PWA", "Next.js"),
        ]),
        ("EDGE + API", BLUE, [
            ("API Gateway", "Amazon"),
            ("CloudFront", "CDN + WAF"),
        ]),
        ("COMPUTE", GREEN, [
            ("Fargate", "Django API"),
            ("Lambda", "LR inference"),
        ]),
        ("AI / ML", NAVY, [
            ("Bedrock", "Claude · vernacular"),
            ("SageMaker", "LR retraining"),
        ]),
        ("DATA", GOLD, [
            ("RDS Postgres", "Multi-AZ"),
            ("DocumentDB", "Score events"),
        ]),
    ]

    lh = (CONTENT_H - Inches(0.9)) / len(layers)
    for i, (title, accent, aws_items) in enumerate(layers):
        y = top + i * (lh + Inches(0.05))
        rect(slide, main_left, y, main_w, lh, WHITE, LINE, radius=True)
        rect(slide, main_left, y, Inches(0.12), lh, accent)
        txt(slide, main_left + Inches(0.25), y + Inches(0.05),
            Inches(1.5), Inches(0.22),
            title, size=8, bold=True, color=NAVY)

        # AWS service badges in a row
        badge_area_x = main_left + Inches(1.85)
        badge_area_w = main_w - Inches(1.95)
        bw = (badge_area_w - Inches(0.1)) / len(aws_items)
        for j, (service, detail) in enumerate(aws_items):
            bx = badge_area_x + j * (bw + Inches(0.05))
            by = y + Inches(0.06)
            aws_badge(slide, bx, by, bw - Inches(0.05), lh - Inches(0.12),
                      service, detail)

        if i < len(layers) - 1:
            arrow(slide, main_left + main_w / 2, y + lh,
                  main_left + main_w / 2, y + lh + Inches(0.05),
                  NAVY_SOFT, 1)

    # Regulated rails band at the bottom of the left column
    rails_top = top + len(layers) * (lh + Inches(0.05))
    rails_h = Inches(0.42)
    rect(slide, main_left, rails_top, main_w, rails_h, NAVY, radius=True)
    rect(slide, main_left, rails_top, Inches(0.12), rails_h, LIME)
    txt(slide, main_left + Inches(0.25), rails_top + Inches(0.09),
        Inches(1.7), Inches(0.24),
        "REGULATED RAILS", size=8, bold=True, color=LIME)
    rails = ["AA · Sahamati · Finvu", "ULI · RBIH · 64 lenders",
             "OCEN 4.0 · iSPIRT"]
    rw = (main_w - Inches(2.2)) / 3
    for i, r in enumerate(rails):
        rx = main_left + Inches(2.1) + i * rw
        txt(slide, rx, rails_top + Inches(0.09), rw - Inches(0.1),
            Inches(0.24), r, size=8, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

    # ─── RIGHT column: AWS foundation + external ─────────
    ext_left = main_left + main_w + Inches(0.2)
    ext_w = CONTENT_RIGHT - ext_left
    section_label(slide, ext_left, top - Inches(0.02), ext_w,
                  "AWS Foundation", color=AWS_ORANGE)

    aws_foundation = [
        ("S3", "Consent artifacts · KYC docs"),
        ("ElastiCache", "Redis · session + cache"),
        ("Secrets Manager", "API keys · mTLS certs"),
        ("Route 53", "DNS · udyamai.credit"),
        ("CloudWatch", "Logs · metrics · alarms"),
        ("IAM · KMS", "Access · encryption CMEK"),
        ("CodePipeline", "CI/CD · blue-green"),
    ]
    yy = top + Inches(0.28)
    eh = Inches(0.42)
    for name, sub in aws_foundation:
        aws_badge(slide, ext_left, yy, ext_w, eh, name, sub)
        yy += eh + Inches(0.03)


# ─── SLIDE 8 · Technologies ────────────────────────────────────────────

def slide_technologies(slide):
    headers = ["Layer", "Choice", "Why"]
    rows = [
        ["Frontend", "Next.js 14 · React Native · Tailwind",
         "Edge-fast web + native mobile · one design system"],
        ["Compute", "AWS Fargate + Lambda",
         "Serverless containers · pay-per-request"],
        ["LLM", "AWS Bedrock (Claude) + Gemini 1.5 Flash",
         "Bedrock for prod · Gemini for hackathon speed"],
        ["ML", "AWS SageMaker · custom LR",
         "Auditable · RBI model-risk approvable"],
        ["Database", "AWS RDS Postgres + DocumentDB",
         "Multi-AZ HA · ap-south-1 (Mumbai)"],
        ["Cache · Queue", "AWS ElastiCache Redis",
         "Session + Celery broker"],
        ["Data rails", "AA · GSTN FIP · EPFO · UPI",
         "Regulated consented only, no scraping"],
        ["Cross-lender", "ULI (RBIH) · 64 lenders",
         "Official multi-lender rail"],
        ["Loan app", "OCEN 4.0 · iSPIRT",
         "Open standard, no lock-in"],
        ["CDN · WAF", "CloudFront + AWS WAF + Shield",
         "Global edge · DDoS shield built in"],
        ["CI/CD", "CodePipeline + CodeBuild",
         "Blue-green deploy · rollback in seconds"],
        ["Obs · Sec", "CloudWatch + Secrets Manager + KMS",
         "One control plane · RBI IT Framework aligned"],
    ]
    table(slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.05), CONTENT_W,
          CONTENT_H - Inches(0.1), headers, rows,
          col_widths=[Inches(1.4), Inches(3.6), Inches(4.3)],
          header_size=9, row_size=7.5)


# ─── SLIDE 9 · Cost ────────────────────────────────────────────────────

def slide_cost(slide):
    h1 = ["One-time build", "Cost"]
    r1 = [
        ["AA integration + sandbox cert", "₹1.5 L"],
        ["OCEN LA onboarding", "₹1.0 L"],
        ["Cash-flow LSTM + model risk", "₹0.8 L"],
        ["WhatsApp bot + KYC", "₹0.7 L"],
        ["Design + security audit", "₹1.0 L"],
        ["Total one-time", "₹5.0 L"],
    ]
    tw = (CONTENT_W - Inches(0.25)) / 2
    table(slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.05), tw, Inches(2.5),
          h1, r1, col_widths=[tw - Inches(1.3), Inches(1.3)],
          header_size=10, row_size=9)

    h2 = ["AWS OPEX (monthly)", "Cost / mo"]
    r2 = [
        ["Fargate + Lambda compute", "₹12k"],
        ["RDS Postgres + DocumentDB", "₹15k"],
        ["Bedrock inference + Gemini", "₹18k"],
        ["S3 + CloudFront + WAF", "₹4k"],
        ["AA data pulls · MSG91", "₹15k"],
        ["Total AWS monthly", "₹64k"],
    ]
    table(slide, CONTENT_LEFT + tw + Inches(0.25), CONTENT_TOP + Inches(0.05),
          tw, Inches(2.15), h2, r2,
          col_widths=[tw - Inches(1.3), Inches(1.3)],
          header_fill=GREEN, header_size=10, row_size=9)

    py = CONTENT_TOP + Inches(2.75)
    pw = Inches(3.0)
    pg = Inches(0.15)
    tot = pw * 3 + pg * 2
    px = CONTENT_LEFT + (CONTENT_W - tot) / 2
    stat_pill(slide, px, py, "Per MSME / month", "< ₹1", SAFFRON, pw,
              Inches(0.7), value_size=20)
    stat_pill(slide, px + pw + pg, py, "Break-even leads / mo", "100", GREEN,
              pw, Inches(0.7), value_size=20)
    stat_pill(slide, px + 2 * (pw + pg), py, "Revenue at 50k MAU", "₹5 L / mo",
              BLUE, pw, Inches(0.7), value_size=20)

    callout(slide, CONTENT_LEFT, CONTENT_BOTTOM - Inches(0.55), CONTENT_W,
            "< ₹1 per MSME · Break-even: 1 bank × ₹500/converted-lead × 100 leads/mo",
            accent=SAFFRON)


# ─── SLIDE 11 · Performance ────────────────────────────────────────────

def slide_performance(slide):
    h1 = ["Lender", "Acc", "AUC", "Bias", "Top β"]
    r1 = [
        ["IDBI Bank", "87.3%", "0.729", "-0.14", "Revenue (+1.14)"],
        ["SBI", "82.8%", "0.789", "+0.00", "Compliance (+1.21)"],
        ["HDFC Bank", "88.8%", "0.721", "-0.04", "Revenue (+0.99)"],
    ]
    tw = Inches(5.6)
    table(slide, CONTENT_LEFT, CONTENT_TOP + Inches(0.05), tw, Inches(1.55),
          h1, r1,
          col_widths=[Inches(1.6), Inches(0.8), Inches(0.8), Inches(0.8),
                      Inches(1.6)],
          header_size=10, row_size=9)

    h2 = ["Path", "Cold", "Warm"]
    r2 = [
        ["Landing /", "190 ms", "45 ms"],
        ["/dashboard pipeline", "900 ms", "220 ms"],
        ["/api/health-card", "340 ms", "60 ms"],
        ["/api/ocen/apply", "210 ms", "45 ms"],
        ["Gemini explanation", "1.4 s", "1.1 s"],
    ]
    rlx = CONTENT_LEFT + tw + Inches(0.2)
    rlw = CONTENT_RIGHT - rlx
    table(slide, rlx, CONTENT_TOP + Inches(0.05), rlw, Inches(2.15),
          h2, r2, col_widths=[Inches(1.8), Inches(0.8), Inches(0.8)],
          header_fill=GREEN, header_size=10, row_size=9)

    py = CONTENT_TOP + Inches(1.85)
    stats = [("Track-3 clauses", "10 / 10", SAFFRON),
             ("Languages", "EN · HI · TE", GREEN),
             ("ULI lenders", "64", BLUE)]
    sw = Inches(1.8)
    sg = Inches(0.15)
    for i, (label, val, color) in enumerate(stats):
        x = CONTENT_LEFT + i * (sw + sg)
        stat_pill(slide, x, py, label, val, color, sw, Inches(0.7),
                  value_size=15)

    callout(slide, CONTENT_LEFT, CONTENT_BOTTOM - Inches(0.55), CONTENT_W,
            "AUC 0.72-0.79 across all 3 lenders · every β visible in the in-product Model Card")


# ─── SLIDE 12 · Roadmap ────────────────────────────────────────────────

def slide_roadmap(slide):
    horizons = [
        ("0-6 MONTHS", SAFFRON, [
            "Real AA integration (Finvu → prod)",
            "UPI counterparty graph analysis",
            "90-day cash-flow LSTM",
            "Real OCEN LA registration",
        ]),
        ("6-18 MONTHS", GREEN, [
            "Khatabook / OkCredit embed (30M+ MSMEs)",
            "WhatsApp vernacular coaching bot",
            "Sector cohort benchmarking",
            "UdyamAI Current Account",
        ]),
        ("18-36 MONTHS", NAVY, [
            "RBI RegTech sandbox partner",
            "Insurance rail on same AA data",
            "International · Indonesia · Vietnam",
            "Series A ready · 500k MSMEs",
        ]),
    ]
    cw = (CONTENT_W - Inches(0.3)) / 3
    ch = CONTENT_H - Inches(0.1)
    for i, (title, accent, items) in enumerate(horizons):
        x = CONTENT_LEFT + i * (cw + Inches(0.15))
        y = CONTENT_TOP + Inches(0.05)
        rect(slide, x, y, cw, ch, WHITE, LINE, radius=True)
        rect(slide, x, y, cw, Inches(0.4), accent, radius=True)
        tc = TEXT if accent == LIME else WHITE
        txt(slide, x + Inches(0.15), y + Inches(0.06), cw - Inches(0.2),
            Inches(0.28), title, size=11, bold=True, color=tc)
        iy = y + Inches(0.55)
        for item in items:
            oval(slide, x + Inches(0.18), iy + Inches(0.12), Inches(0.09),
                 accent)
            txt(slide, x + Inches(0.36), iy, cw - Inches(0.42), Inches(0.45),
                item, size=9, color=INK, line_spacing=1.15)
            iy += Inches(0.55)


# ─── SLIDE 13 · Links ──────────────────────────────────────────────────

def slide_links(slide):
    links = [
        ("GitHub Public Repository",
         "https://github.com/<your-handle>/udyamai", SAFFRON),
        ("Demo Video (3 minutes)",
         "https://youtu.be/<video-id>", GREEN),
        ("Final Product Link",
         "https://udyamai.vercel.app", BLUE),
    ]
    y = CONTENT_TOP + Inches(0.15)
    lh = Inches(1.05)
    gap = Inches(0.15)
    for title, url, accent in links:
        rect(slide, CONTENT_LEFT, y, CONTENT_W, lh, WHITE, LINE, radius=True)
        rect(slide, CONTENT_LEFT, y, Inches(0.12), lh, accent)
        txt(slide, CONTENT_LEFT + Inches(0.3), y + Inches(0.2),
            Inches(6), Inches(0.32), title, size=13, bold=True, color=NAVY)
        txt(slide, CONTENT_LEFT + Inches(0.3), y + Inches(0.58),
            CONTENT_W - Inches(0.5), Inches(0.3), url, size=11, color=BLUE,
            font=MONO)
        y += lh + gap


# ─── driver ────────────────────────────────────────────────────────────

def build():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    slide_brief_idea(slides[1])
    slide_opportunities(slides[2])
    slide_features(slides[3])
    slide_flowchart(slides[4])
    slide_screenshots(slides[5])
    slide_architecture(slides[6])
    slide_technologies(slides[7])
    slide_cost(slides[8])
    slide_screenshots(slides[9])
    slide_performance(slides[10])
    slide_roadmap(slides[11])
    slide_links(slides[12])

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
